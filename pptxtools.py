import re
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

color_dict = {
    'black': RGBColor(0, 0, 0),
    'white': RGBColor(255, 255, 255),
    'red': RGBColor(255, 0, 0),
    'green': RGBColor(0, 255, 0),
    'blue': RGBColor(0, 0, 255),
}


def add_slide(prs, layout):
    return prs.slides.add_slide(layout)


def customizer_topics(slide, placeholder_number: int, font_name, font_size: int, bold: bool, color_name):
    slide.placeholders[placeholder_number].text_frame.paragraphs[0].font.name = font_name
    slide.placeholders[placeholder_number].text_frame.paragraphs[0].font.size = Pt(
        font_size)
    slide.placeholders[placeholder_number].text_frame.paragraphs[0].font.bold = bold
    slide.placeholders[placeholder_number].text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER #just now
    for run in slide.placeholders[placeholder_number].text_frame.paragraphs[0].runs:
        if color_name.lower() in color_dict:
            color = color_dict[color_name.lower()]
            run.font.color.rgb = color
        else:
            run.font.color.rgb = RGBColor(0, 0, 0)


def customizer_bullet_point(slide, placeholder_number: int, paragraph_number: int, font_name, font_size: int, bold: bool, color_name):
    slide.shapes.placeholders[placeholder_number].text_frame.paragraphs[paragraph_number].font.name = font_name
    slide.shapes.placeholders[placeholder_number].text_frame.paragraphs[paragraph_number].font.size = Pt(
        font_size)
    slide.shapes.placeholders[placeholder_number].text_frame.paragraphs[paragraph_number].font.bold = bold
    for run in slide.shapes.placeholders[placeholder_number].text_frame.paragraphs[paragraph_number].runs:
        if color_name.lower() in color_dict:
            color = color_dict[color_name.lower()]
            run.font.color.rgb = color
        else:
            run.font.color.rgb = RGBColor(0, 0, 0)


def customizer_background_color(prs, color):

    # Access the slide master
    slide_master = prs.slide_master

    # Access the background of the slide master
    background = slide_master.background
    background.fill.solid()

    # Set the background color
    if color.lower() in color_dict:
        color = color_dict[color.lower()]
        background.fill.fore_color.rgb = color
    else:
        background.fill.fore_color.rgb = color_dict['white']


def split_sentences(text):
    # sentences = re.split(
    #     r'(?<!\d\.\d.)(?<!\d\.\d)(?<![A-Z][a-z]\.)(?<!\w\.\w.)(?<!\.\.)(?<=\.|\?)\s', text)
    sentences = text.split('\n')
    return sentences


def shape_set_font_size(shape, points):
    for paragraph in shape.text_frame.paragraphs:
        paragraph.font.size = Pt(points)
