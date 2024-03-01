import pdftitle
import pdfplumber
import path
import pptx

from bs4 import BeautifulSoup
from fastapi.middleware.cors import CORSMiddleware
from fastapi import FastAPI, UploadFile, HTTPException, Request
from classobjects import PDF, PresentationData,ImgCap
from pdftools import *
from pptxtools import *
from gemini import gemini_summarize
from presentify_model import summarize
from image_extraction import image_extraction
from cosine_similarity import cosine_similarity
from create_presentation import create_presentation,create_presentation_themechange
from image_slides import display_slides
from pydantic import BaseModel

pdf = PDF()
presentation = PresentationData()
presentation1 = PresentationData()
imgcap = ImgCap()

app = FastAPI()

MAX_FILE_SIZE = 1024 * 1024 * 12

app.add_middleware(
    CORSMiddleware, 
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Define a Pydantic model for the incoming JSON data
class ThemeSelectData(BaseModel):
    theme: str
    
@app.post('/extract-text')
async def extract_texts(file: UploadFile):
    if file.size > MAX_FILE_SIZE:
        raise HTTPException(status_code=413, detail="File is too large")
    pdf_bytes = await file.read()
    print("Starting Website.....")
    # Optionally save the PDF temporarily
    with open('temp_pdf.pdf', 'wb') as f:
        f.write(pdf_bytes)
    pdf.textdata = clean_text(read_pdf('temp_pdf.pdf'))
    try:
        pdf_title = pdftitle.get_title_from_file('temp_pdf.pdf') #ya ho error aayeko
    except:
        pdf_title = 'title'
    try:
        author = pdfplumber.open('temp_pdf.pdf').metadata['Author']
    except:    
        author = 'author name'
        
    presentation.author = author
    presentation.title = pdf_title
    gemini_data = PresentationData()

    try:
        gemini_data = gemini_summarize(pdf.textdata)
        presentation.introduction = clean_text(gemini_data.introduction)
        presentation.literature_review = clean_text(gemini_data.literature_review)
        presentation.methodology = clean_text(gemini_data.methodology)
        presentation.results = clean_text(gemini_data.results)
        presentation.conclusions = clean_text(gemini_data.conclusions)
    except:
        return {'error': 'couldnt extract data'}
    data = PresentationData()
    data = summarize(gemini_data)
    presentation1 = data
    data_dict = {'Introduction': data.introduction,
                 'Literature Review': data.literature_review,
                 'Methodology': data.methodology,
                 'Results': data.results,
                 'Conclusion': data.conclusions}
    image_caption_list = image_extraction(path.temp_path)
    if image_caption_list != []:
        similarity = cosine_similarity(data_dict, image_caption_list)
        filtered_similarity = [item for item in similarity if all(value > 0.25 for value in item.values())]
    else:
        filtered_similarity =[]
    imgcap.imgcap = filtered_similarity
    theme_select_path = r'theme\default.pptx'    
    prs1 = pptx.Presentation(theme_select_path)
    create_presentation(prs1,data_dict, presentation.title,presentation.author, filtered_similarity)
    display_slides()
    return {"message": "Default Slide created successfully!"}

@app.post('/theme-select')
async def theme_select(theme_data: ThemeSelectData):
    selected_theme = theme_data.theme
    theme_select_path = 'theme'+ rf'\{selected_theme}.pptx'
    prs2 = pptx.Presentation(theme_select_path)
    data = presentation1
    data_dict = {'Introduction': data.introduction,
                 'Literature Review': data.literature_review,
                 'Methodology': data.methodology,
                 'Results': data.results,
                 'Conclusion': data.conclusions}
    create_presentation_themechange(prs2,data_dict, presentation.title,presentation.author,  imgcap.imgcap)
    display_slides()
    return {"message": "Slide Created in: {}".format(selected_theme)}

@app.post("/get_data_from_url")
async def get_data_fromI_url(url: str):
    if 'https://arxiv.org/abs' in url:
        response = requests.get(url)
        response = response.content
        soup = BeautifulSoup(response, 'html.parser')

        presentation.title =  soup.find('h1', class_='title mathjax').text
        presentation.author = soup.find('div', class_='authors').text
        pdf_link = url.replace('abs', 'pdf')
        presentation.title = presentation.title.replace('Title:','')
        presentation.author = presentation.author.replace('Authors:','')
        pdf.textdata = clean_text(read_pdf_from_url(pdf_link))
    else:
        pdf_link = url
        pdf.textdata = clean_text(read_pdf_from_url(pdf_link))
        presentation.title = pdftitle.get_title_from_file('temp_pdf.pdf')
        presentation.author = pdfplumber.open('temp_pdf.pdf').metadata['Author']
    gemini_data = PresentationData()
    # gemini_data = gemini_summarize(pdf.textdata)
    try:
        gemini_data = gemini_summarize(pdf.textdata)
        presentation.introduction = gemini_data.introduction
        presentation.literature_review = gemini_data.literature_review
        presentation.methodology = gemini_data.methodology
        presentation.results = gemini_data.results
        presentation.conclusions = gemini_data.conclusions
    except:
        return {'error':'couldnt extract data'}
    data = PresentationData()
    data = summarize(gemini_data)
    data_dict = {'Introduction': data.introduction,
                 'Literature Review': data.literature_review,
                 'Methodology': data.methodology,
                 'Results': data.results,
                 'Conclusion': data.conclusions}
    image_caption_list = image_extraction(path.temp_path)
    if image_caption_list != []:
        similarity = cosine_similarity(data_dict, image_caption_list)
        filtered_similarity = [item for item in similarity if all(value > 0.25 for value in item.values())]
    else:
        filtered_similarity =[]
    theme_select_path = r'theme\default.pptx'    
    prs = pptx.Presentation(theme_select_path)
    create_presentation(prs,data_dict, presentation.title,presentation.author, filtered_similarity)
    display_slides()
    return {"message": "Slide created successfully!"}

