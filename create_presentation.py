import pptx
from pptx.util import Inches
from classobjects import PresentationData

from pptxtools import *
from quilbot import bullet_point

presentation2 = PresentationData()
title_list = ['Introduction', 'Literature Review',
                  'Methodology', 'Results', 'Conclusion']

title_color = "black"
author_color = "black"
bullet_color = "black"
bullet_title_color = "black"
background_color = "white"
def remove_extra_spaces(text):
    return ' '.join(text.split())

presentation_filename = 'empty.pptx'

def create_presentation(prs,data_dict, presentation_title, presentation_author, filtered_similarity):

    # path = r'theme\default.pptx'    
    # prs = pptx.Presentation(path)

    # prs.slide_width = Inches(16)
    # prs.slide_height = Inches(9)

    # title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]
    image_slide_layout = prs.slide_layouts[5]

    # addding page 1
    # slide1 = add_slide(prs, title_slide_layout)
    slide1 = prs.slides[0]
    
    slide1.placeholders[0].text = presentation_title
    # slide1.placeholders[0].height = Inches(1.8)
    # slide1.placeholders[0].width = Inches(9.94)
    # slide1.placeholders[0].left = Inches(2.63)
    # slide1.placeholders[0].top = Inches(1.65)
    

    slide1.placeholders[1].text = presentation_author
    # slide1.placeholders[1].height = Inches(1.2)
    # slide1.placeholders[1].width = Inches(9.94)
    # slide1.placeholders[1].left = Inches(4.43)
    # slide1.placeholders[1].top = Inches(1.65)

    # font style for title
    customizer_topics(slide1, 0, 'Arial', 44, True, title_color)
    # font style for author
    customizer_topics(slide1, 1, 'Arial', 22, True, author_color)
    # filtered_similarity
    # bullet slides
    for title in title_list:
        slide2 = add_slide(prs, bullet_slide_layout)
        slide2.placeholders[0].text = title
        
        slide2.placeholders[0].height = Inches(1.44)
        slide2.placeholders[0].width = Inches(12.36)
        slide2.placeholders[0].left = Inches(0.74)
        slide2.placeholders[0].top = Inches(0.67)

        content = data_dict[title]
        content = bullet_point(content) #theme wala file ma pathau
        sentences = split_sentences(content)
        data_dict[title] = content
        slide2.shapes.placeholders[1].text_frame.text = sentences[0]

        slide2.placeholders[1].height = Inches(4.84)
        slide2.placeholders[1].width = Inches(12.36)
        slide2.placeholders[1].left = Inches(0.74)
        slide2.placeholders[1].top = Inches(2.36)

        customizer_bullet_point(slide2, 1, 0, 'Arial', 28, False, bullet_color)

        for i in range(1, len(sentences)):
            p = slide2.shapes.placeholders[1].text_frame.add_paragraph()
            p.text = sentences[i]
            p.level = 0
            # font style for bullet
            customizer_bullet_point(slide2, 1, i, 'Arial', 28, False, bullet_color)

        # font style for title
        customizer_topics(slide2, 0, 'Arial', 36, True, bullet_title_color)
        for item in filtered_similarity:
            if title in item:
                image_index = item['Index']
                img_path = "images/"+str(image_index)+".jpg"
                slide3 = add_slide(prs, image_slide_layout)
                slide3.placeholders[0].text = title
                slide2.placeholders[0].height = Inches(1.44)
                slide2.placeholders[0].width = Inches(12.36)
                slide2.placeholders[0].left = Inches(0.74)
                slide2.placeholders[0].top = Inches(0.67)
                
                customizer_topics(slide3, 0, 'Arial', 44, True, title_color)
                
                height = Inches(5)
                width = Inches(8)
                left = Inches(1.3)
                top = Inches(1.6)
                
                img = slide3.shapes.add_picture(img_path, 0,0,height = height,width = width)
        
                img.left = int(left)
                img.top = int(top)

    customizer_background_color(prs, background_color)
    presentation_filename = ''
    for char in presentation_title:
        if char.isalnum() or char == ' ':
            presentation_filename += char
        else:
            presentation_filename += ' '
    presentation_filename = remove_extra_spaces(presentation_filename)
    prs.save(f'slides/{presentation_filename}.pptx')
    presentation2.title = presentation_title
    presentation2.author = presentation_author
    presentation2.introduction = data_dict['Introduction']
    presentation2.literature_review = data_dict['Literature Review']
    presentation2.methodology = data_dict['Methodology']
    presentation2.results = data_dict['Results']
    presentation2.conclusions = data_dict['Conclusion']

def create_presentation_themechange(prs,data_dict, presentation_title, presentation_author, filtered_similarity):

    bullet_slide_layout = prs.slide_layouts[1]
    image_slide_layout = prs.slide_layouts[5]

    # addding page 1
    # slide1 = add_slide(prs, title_slide_layout)
    slide1 = prs.slides[0]
    
    slide1.placeholders[0].text = presentation_title
    # slide1.placeholders[0].height = Inches(1.8)
    # slide1.placeholders[0].width = Inches(9.94)
    # slide1.placeholders[0].left = Inches(2.63)
    # slide1.placeholders[0].top = Inches(1.65)
    

    slide1.placeholders[1].text = presentation_author
    # slide1.placeholders[1].height = Inches(1.2)
    # slide1.placeholders[1].width = Inches(9.94)
    # slide1.placeholders[1].left = Inches(4.43)
    # slide1.placeholders[1].top = Inches(1.65)

    # font style for title
    customizer_topics(slide1, 0, 'Arial', 44, True, title_color)
    # font style for author
    customizer_topics(slide1, 1, 'Arial', 22, True, author_color)
    # filtered_similarity
    # bullet slides
    presentation2.title = presentation_title
    presentation2.author = presentation_author
    data_dict['Introduction'] = presentation2.introduction
    data_dict['Literature Review'] = presentation2.literature_review
    data_dict['Methodology'] = presentation2.methodology 
    data_dict['Results'] = presentation2.results
    data_dict['Conclusion'] = presentation2.conclusions  
    
    for title in title_list:
        slide2 = add_slide(prs, bullet_slide_layout)
        slide2.placeholders[0].text = title
        
        slide2.placeholders[0].height = Inches(1.44)
        slide2.placeholders[0].width = Inches(12.36)
        slide2.placeholders[0].left = Inches(0.74)
        slide2.placeholders[0].top = Inches(0.67)
        
        content = data_dict[title]
        sentences = split_sentences(content)
        slide2.shapes.placeholders[1].text_frame.text = sentences[0]

        slide2.placeholders[1].height = Inches(4.84)
        slide2.placeholders[1].width = Inches(12.36)
        slide2.placeholders[1].left = Inches(0.74)
        slide2.placeholders[1].top = Inches(2.36)

        customizer_bullet_point(slide2, 1, 0, 'Arial', 28, False, bullet_color)

        for i in range(1, len(sentences)):
            p = slide2.shapes.placeholders[1].text_frame.add_paragraph()
            p.text = sentences[i]
            p.level = 0
            # font style for bullet
            customizer_bullet_point(slide2, 1, i, 'Arial', 28, False, bullet_color)

        # font style for title
        customizer_topics(slide2, 0, 'Arial', 36, True, bullet_title_color)
        for item in filtered_similarity:
            if title in item:
                image_index = item['Index']
                img_path = "images/"+str(image_index)+".jpg"
                slide3 = add_slide(prs, image_slide_layout)
                slide3.placeholders[0].text = title
                slide2.placeholders[0].height = Inches(1.44)
                slide2.placeholders[0].width = Inches(12.36)
                slide2.placeholders[0].left = Inches(0.74)
                slide2.placeholders[0].top = Inches(0.67)
                
                customizer_topics(slide3, 0, 'Arial', 44, True, title_color)
                
                height = Inches(5)
                width = Inches(8)
                left = Inches(1.3)
                top = Inches(1.6)
                
                img = slide3.shapes.add_picture(img_path, 0,0,height = height,width = width)
        
                img.left = int(left)
                img.top = int(top)

    customizer_background_color(prs, background_color)
    presentation_filename = ''
    for char in presentation_title:
        if char.isalnum() or char == ' ':
            presentation_filename += char
        else:
            presentation_filename += ' '
    presentation_filename = remove_extra_spaces(presentation_filename)
    prs.save(f'slides/{presentation_filename}.pptx')
